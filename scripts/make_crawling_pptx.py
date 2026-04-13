from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from textwrap import wrap

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
README_PATH = ROOT / "README.md"
OUT_PATH = ROOT / "크롤링방법정리.pptx"


@dataclass(frozen=True)
class Theme:
    bg: RGBColor = RGBColor(0xF2, 0xF1, 0xED)
    fg: RGBColor = RGBColor(0x26, 0x25, 0x1E)
    muted: RGBColor = RGBColor(0x6B, 0x67, 0x5D)  # warm-ish muted
    accent: RGBColor = RGBColor(0xF5, 0x4E, 0x00)
    error: RGBColor = RGBColor(0xCF, 0x2D, 0x56)
    card: RGBColor = RGBColor(0xE6, 0xE5, 0xE0)


THEME = Theme()


def _set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = THEME.bg


def _add_title(slide, text: str) -> None:
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(12.0), Inches(0.9))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(34)
    run.font.bold = False
    run.font.color.rgb = THEME.fg
    p.alignment = PP_ALIGN.LEFT


def _add_subtitle(slide, text: str) -> None:
    tb = slide.shapes.add_textbox(Inches(0.82), Inches(1.28), Inches(12.0), Inches(0.55))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(16)
    run.font.color.rgb = THEME.muted
    p.alignment = PP_ALIGN.LEFT


def _add_card(slide, x, y, w, h, title: str, bullets: list[str]) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = THEME.card
    shape.line.color.rgb = RGBColor(0xC9, 0xC5, 0xBC)
    shape.line.width = Pt(1.0)

    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.18)
    tf.word_wrap = True

    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = title
    r0.font.size = Pt(18)
    r0.font.color.rgb = THEME.fg
    r0.font.bold = False

    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(14)
        p.font.color.rgb = THEME.fg


def _add_code_block(slide, x, y, w, h, code: str) -> None:
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0xF7, 0xF7, 0xF4)
    box.line.color.rgb = RGBColor(0xC9, 0xC5, 0xBC)
    box.line.width = Pt(1.0)

    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.18)
    tf.word_wrap = True

    lines = code.strip("\n").splitlines()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line.rstrip()
        p.level = 0
        p.font.size = Pt(12)
        p.font.color.rgb = THEME.fg
        p.font.name = "Consolas"


def _add_flow(slide) -> None:
    # Simple flow diagram using shapes
    def node(x, y, w, h, text, fill):
        s = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, x, y, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = fill
        s.line.color.rgb = RGBColor(0xC9, 0xC5, 0xBC)
        s.line.width = Pt(1.0)
        tf = s.text_frame
        tf.clear()
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.size = Pt(14)
        r.font.color.rgb = THEME.fg
        return s

    n1 = node(Inches(0.9), Inches(2.2), Inches(3.2), Inches(1.0), "HITL 로그인\\n(Chrome + hCaptcha)", RGBColor(0xEB, 0xEA, 0xE5))
    n2 = node(Inches(4.4), Inches(2.2), Inches(3.2), Inches(1.0), "크롤링 실행\\n(main.py / Streamlit)", RGBColor(0xE6, 0xE5, 0xE0))
    n3 = node(Inches(7.9), Inches(2.2), Inches(3.2), Inches(1.0), "정규화(JSON) 저장\\nresult/<YYMM>/*.json", RGBColor(0xEB, 0xEA, 0xE5))
    n4 = node(Inches(4.4), Inches(3.9), Inches(3.2), Inches(1.0), "대시보드\\n(Next.js)", RGBColor(0xE1, 0xE0, 0xDB))

    # arrows (simple lines)
    def arrow(x1, y1, x2, y2):
        line = slide.shapes.add_connector(1, x1, y1, x2, y2)
        line.line.color.rgb = RGBColor(0xB0, 0xAE, 0xA5)
        line.line.width = Pt(2.0)
        line.line.end_arrowhead = True

    arrow(n1.left + n1.width, n1.top + n1.height / 2, n2.left, n2.top + n2.height / 2)
    arrow(n2.left + n2.width, n2.top + n2.height / 2, n3.left, n3.top + n3.height / 2)
    arrow(n3.left + n3.width / 2, n3.top + n3.height, n4.left + n4.width / 2, n4.top)


def _readme_snippets(readme: str) -> dict[str, str]:
    # Keep minimal: pull key command blocks if present.
    blocks: list[str] = []
    cur: list[str] = []
    in_code = False
    for line in readme.splitlines():
        if line.strip().startswith("```"):
            if in_code:
                blocks.append("\n".join(cur))
                cur = []
                in_code = False
            else:
                in_code = True
            continue
        if in_code:
            cur.append(line.rstrip("\n"))
    # pick first few meaningful blocks
    def pick_contains(substr: str) -> str | None:
        for b in blocks:
            if substr in b:
                return b
        return None

    return {
        "install": pick_contains("pip install") or "",
        "run_cli": pick_contains("python main.py") or "",
        "run_streamlit": pick_contains("streamlit run streamlit_app.py") or "",
        "run_dashboard": pick_contains("npm run dev") or "",
        "batch": pick_contains("--batch-all-sections") or "",
    }


def build_pptx() -> None:
    readme = README_PATH.read_text(encoding="utf-8")
    snips = _readme_snippets(readme)

    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # Slide 1: Title
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title(s, "eJARN 크롤링 방법 정리")
    _add_subtitle(s, f"README 요약 기반 · 생성일 {datetime.now().strftime('%Y-%m-%d')}")
    _add_card(
        s,
        Inches(0.8),
        Inches(2.1),
        Inches(12.0),
        Inches(4.7),
        "이 문서가 설명하는 것",
        [
            "HITL(사람) 로그인 기반 수집 흐름",
            "결과 JSON 스키마와 저장 규칙",
            "배치 수집(8개 섹션) 실행 방법",
            "로컬 대시보드(월별 2603/2604 요약) 실행 방법",
        ],
    )

    # Slide 2: End-to-end flow
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title(s, "전체 흐름(End-to-End)")
    _add_subtitle(s, "로그인 → 크롤링 → 정규화(JSON) → 분석/대시보드")
    _add_flow(s)

    # Slide 3: Requirements & env
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title(s, "준비사항(환경/설치/.env)")
    _add_subtitle(s, "Python/Playwright 설치 + .env 계정 설정")
    left = [
        "Python 3.10+ (권장 3.12)",
        "Playwright + Chromium 설치",
        ".env: EJARN_LOGIN_EMAIL / EJARN_LOGIN_PASSWORD 필수",
        "선택: OPENAI_API_KEY (요약/분류/챗봇)",
    ]
    _add_card(s, Inches(0.8), Inches(2.0), Inches(6.1), Inches(4.9), "핵심 체크리스트", left)
    if snips["install"].strip():
        _add_code_block(s, Inches(7.1), Inches(2.0), Inches(5.7), Inches(2.4), snips["install"])
    _add_code_block(
        s,
        Inches(7.1),
        Inches(4.6),
        Inches(5.7),
        Inches(2.3),
        "copy .env.sample .env\n# 그리고 .env에 로그인 계정 입력\nEJARN_LOGIN_EMAIL=...\nEJARN_LOGIN_PASSWORD=...",
    )

    # Slide 4: Crawl (CLI / Streamlit)
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title(s, "크롤링 실행(기본)")
    _add_subtitle(s, "CLI 또는 Streamlit(HITL)로 실행")
    _add_card(
        s,
        Inches(0.8),
        Inches(1.95),
        Inches(6.1),
        Inches(5.0),
        "HITL 동작 포인트",
        [
            "Chrome 창에서 hCaptcha/로그인 완료 필요",
            "CLI: 로그인 완료 후 터미널에 '완료' 입력",
            "Streamlit: '✅ 로그인 완료(진행)' 버튼으로 다음 단계",
            "결과는 result/ 폴더에 JSON 저장",
        ],
    )
    code = []
    if snips["run_cli"].strip():
        code.append(snips["run_cli"].strip())
    if snips["run_streamlit"].strip():
        code.append(snips["run_streamlit"].strip())
    if code:
        _add_code_block(s, Inches(7.1), Inches(1.95), Inches(5.7), Inches(5.0), "\n\n".join(code))

    # Slide 5: Batch crawl
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title(s, "전 섹션 배치 수집(월 폴더 저장)")
    _add_subtitle(s, "한 번 로그인 후 8개 섹션을 순회 수집")
    _add_card(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(6.1),
        Inches(4.9),
        "수집 대상(8개)",
        [
            "Jarn Regular / Jarn Special",
            "eJarn News / Cover Story",
            "Event > Exhibition / Report",
            "Special Issue / Regular Issue",
        ],
    )
    if snips["batch"].strip():
        _add_code_block(s, Inches(7.1), Inches(2.0), Inches(5.7), Inches(4.9), snips["batch"])

    # Slide 6: Output JSON schema + quality
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title(s, "결과 JSON 구조(정규화)")
    _add_subtitle(s, "ArticleCollection → items[ArticleItem] (태그 리스트 기반)")
    _add_code_block(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(12.0),
        Inches(3.2),
        """{
  "source": "eJARN",
  "collected_at": "ISO",
  "items": [
    {
      "date": "YYYY-MM-DD",
      "topic": "...",
      "summary": "... (<= 900자)",
      "link": "https://...",
      "company": [],
      "product_type": [],
      "market_segment": [],
      "refrigerant": ["Unknown"|...],
      "application": ["Multi-purpose"|...],
      "technology": [],
      "category": ["Market"|...]
    }
  ]
}""",
    )
    _add_card(
        s,
        Inches(0.8),
        Inches(5.4),
        Inches(12.0),
        Inches(1.5),
        "품질 규칙(핵심)",
        [
            "summary 최대 900자, 분류 리스트 중복 제거",
            "refrigerant 미식별: Unknown / application 미식별: Multi-purpose / category 미식별: Market",
        ],
    )

    # Slide 7: Dashboard
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title(s, "대시보드(로컬, Next.js)")
    _add_subtitle(s, "result/YYMM/*.json을 읽어 월별(2603/2604) 요약과 품질 신호를 제공")
    _add_card(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(6.1),
        Inches(4.9),
        "현재 제공 기능(MVP)",
        [
            "Monthly: 월별 KPI/TopN/월간 브리프(규칙 기반)",
            "Explore: 최근 파일 일부 미리보기(중복 링크 대비)",
            "Quality: 404/기본값 신호 + 파일별 품질",
            "Trends: 월별 추이(확장 예정)",
        ],
    )
    dash_code = snips["run_dashboard"].strip() or "cd dashboard-frontend\nnpm install\nnpm run dev"
    _add_code_block(s, Inches(7.1), Inches(2.0), Inches(5.7), Inches(2.2), dash_code)
    _add_code_block(
        s,
        Inches(7.1),
        Inches(4.5),
        Inches(5.7),
        Inches(2.4),
        "/monthly?yymm=2604\n/explore\n/quality\n/trends",
    )

    # Slide 8: Mermaid diagram (optional)
    s = prs.slides.add_slide(blank)
    _set_slide_bg(s)
    _add_title(s, "UML(mermaid) — 참고")
    _add_subtitle(s, "문서/위키에 붙여 넣어 흐름 공유할 때 사용")
    _add_code_block(
        s,
        Inches(0.8),
        Inches(2.0),
        Inches(12.0),
        Inches(4.8),
        """```mermaid
flowchart LR
  HitlLogin[HITL_Login] --> Crawl[Crawl_Pipeline]
  Crawl --> Normalize[Normalize_JSON]
  Normalize --> ResultDir[result_YYMM_json]
  ResultDir --> Dashboard[Nextjs_Dashboard]
```""",
    )
    _add_card(
        s,
        Inches(0.8),
        Inches(6.95),
        Inches(12.0),
        Inches(0.4),
        "Tip",
        ["README 기반으로 생성됨(세부는 README/코드 참조)."],
    )

    prs.save(OUT_PATH)


if __name__ == "__main__":
    if not README_PATH.exists():
        raise SystemExit(f"README not found: {README_PATH}")
    build_pptx()
    print(f"saved: {OUT_PATH}")

